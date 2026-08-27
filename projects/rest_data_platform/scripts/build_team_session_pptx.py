"""Build ASTRA IMC team-session PowerPoint (Eisberg vor UI)."""

from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from lxml import etree

ROOT = Path(__file__).resolve().parents[1]
MEDIA = ROOT / "04_communication" / "media" / "team_session_2026_08"
OUT = ROOT / "04_communication" / "ASTRA_IMC_Team_Stand_2026_08.pptx"

NAVY = RGBColor(0x0B, 0x2E, 0x3D)
TEAL = RGBColor(0x1A, 0x7A, 0x7A)
SAND = RGBColor(0xF4, 0xF0, 0xE6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5A, 0x65, 0x6B)
ORANGE = RGBColor(0xC4, 0x6B, 0x2B)
PIL_NAVY = (11, 46, 61)
PIL_TEAL = (26, 122, 122)
PIL_SAND = (244, 240, 230)
PIL_WHITE = (255, 255, 255)
PIL_DEEP = (8, 70, 88)


def font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    name = "C:/Windows/Fonts/segoeuib.ttf" if bold else "C:/Windows/Fonts/segoeui.ttf"
    return ImageFont.truetype(name, size)


def set_run_font(run, name="Calibri", size=18, bold=False, color=INK):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    latin = rPr.find(qn("a:latin"))
    if latin is None:
        latin = etree.SubElement(rPr, qn("a:latin"))
    latin.set("typeface", name)


def add_bg(slide, color=SAND):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def bar(slide, prs):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.12)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.fill.background()


def footer(slide, prs, page: str):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.12), Inches(10.5), Inches(0.28))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "ASTRA IMC  ·  WP 5.2 ReST Data Platform  ·  Team-Session Aug 2026"
    set_run_font(run, size=11, color=MUTED)
    num = slide.shapes.add_textbox(Inches(12.2), Inches(7.12), Inches(0.7), Inches(0.28))
    tf2 = num.text_frame
    tf2.clear()
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = page
    set_run_font(r2, size=11, color=MUTED)


def textbox(slide, l, t, w, h, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run_font(run, size=size, bold=bold, color=color)
    return box


def bullets(slide, l, t, w, h, items, size=16):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = "•  " + item
        set_run_font(run, size=size, color=INK)


def card(slide, l, t, w, h, fill=WHITE):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = RGBColor(0xD8, 0xD2, 0xC4)
    sh.adjustments[0] = 0.08
    return sh


def draw_iceberg(path: Path) -> None:
    w, h = 1920, 1080
    im = Image.new("RGB", (w, h), PIL_SAND)
    d = ImageDraw.Draw(im)
    d.rectangle([0, 430, w, h], fill=(186, 214, 214))
    d.rectangle([0, 418, w, 442], fill=PIL_TEAL)
    tip = [(960, 90), (1140, 418), (780, 418)]
    d.polygon(tip, fill=PIL_WHITE, outline=PIL_NAVY)
    body = [(780, 442), (1140, 442), (1380, 1020), (540, 1020)]
    d.polygon(body, fill=PIL_NAVY)
    d.polygon([(820, 560), (1100, 560), (1240, 780), (680, 780)], fill=PIL_TEAL)
    d.polygon([(680, 780), (1240, 780), (1380, 1020), (540, 1020)], fill=PIL_DEEP)
    d.text((790, 220), "Working Board", font=font(36, True), fill=PIL_NAVY)
    d.text((860, 270), "(dünne UI)", font=font(26), fill=PIL_NAVY)
    d.text((820, 470), "Register  ·  Schema  ·  RLS", font=font(28, True), fill=PIL_WHITE)
    d.text((790, 640), "ETL  ·  Matching  ·  Provenance", font=font(28, True), fill=PIL_WHITE)
    d.text((620, 880), "Quellen:  4C   MaStR   Natura   Häfen   ERA5", font=font(26, True), fill=PIL_WHITE)
    d.text((70, 40), "Die Oberfläche ist die Spitze.", font=font(42, True), fill=PIL_NAVY)
    d.text((70, 100), "Die Arbeit sitzt darunter.", font=font(28), fill=PIL_TEAL)
    im.save(path, "PNG")


def draw_kernmodell(path: Path) -> None:
    w, h = 1920, 1080
    im = Image.new("RGB", (w, h), PIL_SAND)
    d = ImageDraw.Draw(im)
    d.text((70, 40), "Kernmodell  —  Park als Nabe", font=font(42, True), fill=PIL_NAVY)
    d.text((70, 100), "37 imc_*-Tabellen  ·  alles hängt an farm_id + source_id", font=font(24), fill=PIL_TEAL)

    def oval(xy, label, sub, fill=PIL_TEAL):
        d.rounded_rectangle(xy, radius=28, fill=fill)
        x0, y0, x1, y1 = xy
        cx, cy = (x0 + x1) / 2, (y0 + y1) / 2
        f1, f2 = font(22, True), font(16)
        bbox = d.textbbox((0, 0), label, font=f1)
        d.text((cx - (bbox[2] - bbox[0]) / 2, cy - 22), label, font=f1, fill=PIL_WHITE)
        bbox2 = d.textbbox((0, 0), sub, font=f2)
        d.text((cx - (bbox2[2] - bbox2[0]) / 2, cy + 8), sub, font=f2, fill=(220, 236, 236))

    nodes = [
        ((80, 220, 430, 360), "Design / Technik", "Foundation, MW, Tiefe"),
        ((1490, 220, 1840, 360), "MaStR-Einheiten", "1.593 Turbinen"),
        ((80, 430, 430, 570), "Milestones", "Inbetriebnahme, EoL"),
        ((1490, 430, 1840, 570), "Häfen", "118 Ports, 677 Links"),
        ((80, 720, 430, 860), "ERA5 Wetter", "3 Parks, 1.858 Tage"),
        ((1490, 720, 1840, 860), "Schutzgebiete", "205 BfN marin"),
        ((760, 820, 1160, 960), "imc_data_sources", "18 Quellen  ·  Provenance"),
    ]
    cx, cy = 960, 560
    for xy, *_ in nodes:
        x0, y0, x1, y1 = xy
        nx, ny = (x0 + x1) / 2, (y0 + y1) / 2
        d.line([(cx, cy), (nx, ny)], fill=(26, 122, 122), width=4)
    for xy, a, b in nodes:
        oval(xy, a, b)
    hub = (760, 430, 1160, 690)
    d.ellipse(hub, fill=PIL_NAVY)
    d.text((838, 510), "imc_wind_farms", font=font(28, True), fill=PIL_WHITE)
    d.text((845, 555), "Nabe  ·  3.606 Parks", font=font(20), fill=(180, 210, 210))
    d.text((880, 590), "DE aktiv: 76", font=font(20), fill=(180, 210, 210))
    im.save(path, "PNG")


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    bar(slide, prs)
    return slide


def build() -> None:
    MEDIA.mkdir(parents=True, exist_ok=True)
    iceberg = MEDIA / "iceberg_working_board.png"
    kern = MEDIA / "kernmodell_nabe.png"
    draw_iceberg(iceberg)
    draw_kernmodell(kern)

    shot = MEDIA / "supabase_schema_visualizer.png"
    cropped = MEDIA / "supabase_schema_visualizer_crop.png"
    if shot.exists():
        im = Image.open(shot)
        w, h = im.size
        left = int(w * 0.18)
        im.crop((left, int(h * 0.10), w, h)).save(cropped)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 title
    s = new_slide(prs)
    textbox(s, 0.7, 1.6, 12, 0.5, "ASTRA IMC  ·  Team-Session", 20, True, TEAL)
    textbox(s, 0.7, 2.1, 12, 1.4, "Zuerst die Arbeit darunter.\nDann das Working Board.", 40, True, NAVY)
    textbox(
        s,
        0.7,
        4.6,
        11,
        1.2,
        "WP 5.2 ReST Data Platform  ·  Pilot WP 2.1  ·  Stand 17.08.2026\nLokales Register, nicht das Endprodukt.",
        18,
        False,
        MUTED,
    )
    footer(s, prs, "1")

    # 2 why
    s = new_slide(prs)
    textbox(s, 0.6, 0.4, 12, 0.6, "Warum nicht mit der Tabelle starten?", 28, True, NAVY)
    card(s, 0.6, 1.3, 12.1, 1.6)
    textbox(
        s,
        0.85,
        1.5,
        11.6,
        1.2,
        "Wer zuerst die UI sieht, sagt: „simple — wo ist LCA, Simulation, der Pass?“\nDie Oberfläche ist absichtlich dünn. Die Substanz ist das Register.",
        18,
        False,
        INK,
    )
    bullets(
        s,
        0.7,
        3.2,
        12,
        3.2,
        [
            "Reihenfolge heute: Eisberg  →  Quellen  →  Schema  →  Zahlen  →  Live-Board  →  Plan",
            "Fehlendes hat einen Owner (Marc / Thomas / Shubham) — das ist der Schnitt, kein vergessenes Ticket.",
            "Plattform liefert Grundlagen. AnyLogic und SimaPro rechnen nicht hier.",
        ],
        18,
    )
    footer(s, prs, "2")

    # 3 frame
    s = new_slide(prs)
    textbox(s, 0.6, 0.4, 12, 0.5, "Rahmen: zwei Produkte, zwei Schichten", 28, True, NAVY)
    card(s, 0.6, 1.2, 6.0, 3.4)
    textbox(s, 0.8, 1.4, 5.6, 0.4, "Produkt A", 16, True, TEAL)
    textbox(s, 0.8, 1.85, 5.6, 1.0, "Decom Capacity Screener", 22, True, NAVY)
    textbox(
        s,
        0.8,
        2.7,
        5.6,
        1.6,
        "Zeitstrahl, Material grob, Hafen, Wetter grob.\nHeute: Datenbausteine. Noch nicht das fertige Produkt.",
        16,
    )
    card(s, 6.9, 1.2, 6.0, 3.4)
    textbox(s, 7.1, 1.4, 5.6, 0.4, "Produkt B", 16, True, TEAL)
    textbox(s, 7.1, 1.85, 5.6, 1.0, "DPP-Prototyp / AAS", 22, True, NAVY)
    textbox(
        s,
        7.1,
        2.7,
        5.6,
        1.6,
        "Feldliste, Submodels, ein Export-PoC.\nPostgres bleibt System of Record.\nAAS ist die reisefähige Form.",
        16,
    )
    card(s, 0.6, 4.8, 12.1, 1.9)
    textbox(
        s,
        0.85,
        5.0,
        11.6,
        1.5,
        "Dual-Track:  DB = Wahrheit für Marc & Thomas (CSV/View).   AAS = Export für Interop (Shubham).\nNicht: AAS als zweites Register. Nicht: Simulation in der Plattform.",
        17,
    )
    footer(s, prs, "3")

    # 4 iceberg
    s = new_slide(prs)
    textbox(s, 0.5, 0.32, 12, 0.4, "Die Arbeit darunter", 26, True, NAVY)
    s.shapes.add_picture(str(iceberg), Inches(0.45), Inches(0.8), Inches(12.4), Inches(6.15))
    footer(s, prs, "4")

    # 5 sources
    s = new_slide(prs)
    textbox(s, 0.5, 0.32, 12, 0.4, "Was eingeflossen ist", 26, True, NAVY)
    src = MEDIA / "quellen_fluss_pictogramme.png"
    s.shapes.add_picture(str(src), Inches(0.35), Inches(0.75), Inches(12.6), Inches(5.35))
    textbox(
        s,
        0.5,
        6.15,
        12.3,
        0.7,
        "4C = Parks/Technik  ·  MaStR = DE-IDs & Einheiten  ·  BfN = Schutzgebiete  ·  Häfen = Lage/Rolle/Distanz  ·  ERA5 = Tagesreihe Wind/Hs",
        14,
        False,
        MUTED,
    )
    footer(s, prs, "5")

    # 6 schema screenshot
    s = new_slide(prs)
    textbox(s, 0.5, 0.32, 12, 0.4, "Echtes Schema  —  Supabase Visualizer (lokal)", 24, True, NAVY)
    pic = cropped if cropped.exists() else shot
    s.shapes.add_picture(str(pic), Inches(0.4), Inches(0.8), Inches(12.5), Inches(5.55))
    textbox(
        s,
        0.5,
        6.4,
        12.3,
        0.45,
        "Ausschnitt: farm_id-Satelliten + source_id. 37 IMC-Tabellen sitzen auf Makerkit-Auth/RLS — deshalb wirkt der Visualizer voll.",
        13,
        False,
        MUTED,
    )
    footer(s, prs, "6")

    # 7 kernmodell
    s = new_slide(prs)
    s.shapes.add_picture(str(kern), Inches(0.25), Inches(0.25), Inches(12.8), Inches(6.85))
    footer(s, prs, "7")

    # 8 numbers
    s = new_slide(prs)
    textbox(s, 0.6, 0.4, 12, 0.5, "Was schon zählbar ist  (lokal, 17.08.2026)", 26, True, NAVY)
    stats = [
        ("76", "DE-Parks aktiv", "ohne Cancelled"),
        ("63", "mit Inbetriebnahme", "MaStR + 4C Dates"),
        ("1.593", "Turbinen-Einheiten", "MaStR, u. a. AV01–12"),
        ("205", "Schutzgebiete", "BfN marin Overlay"),
        ("118", "Häfen", "677 Farm-Links"),
        ("1.858", "ERA5-Tage", "AV + Albatros + Amrumbank"),
    ]
    for i, (n, title, sub) in enumerate(stats):
        col, row = i % 3, i // 3
        l = 0.6 + col * 4.15
        t = 1.2 + row * 2.55
        card(s, l, t, 3.9, 2.3)
        textbox(s, l + 0.2, t + 0.25, 3.5, 0.7, n, 36, True, TEAL)
        textbox(s, l + 0.2, t + 1.0, 3.5, 0.5, title, 16, True, NAVY)
        textbox(s, l + 0.2, t + 1.5, 3.5, 0.5, sub, 13, False, MUTED)
    footer(s, prs, "8")

    # 9 AV
    s = new_slide(prs)
    textbox(s, 0.6, 0.4, 12, 0.5, "Ein Objekt durchgängig: Alpha Ventus", 26, True, NAVY)
    rows = [
        ("4C", "Park, Tripod, Standort, Lifecycle"),
        ("MaStR", "Park-Match + 12 Einheiten AV01–AV12"),
        ("Hafen", "Emden kuratiert als Decom-Rolle"),
        ("Natura", "nächstes Schutzgebiet am Detail"),
        ("ERA5", "680 Tage  2024-01-01 – 2026-08-12"),
        ("Export", "Wetter-CSV: day, wind_ms, hs_m, operable"),
    ]
    for i, (k, v) in enumerate(rows):
        y = 1.15 + i * 0.85
        card(s, 0.6, y, 12.1, 0.75, WHITE)
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(y + 0.15), Inches(1.7), Inches(0.45))
        sh.fill.solid()
        sh.fill.fore_color.rgb = TEAL
        sh.line.fill.background()
        textbox(s, 0.75, y + 0.18, 1.7, 0.4, k, 14, True, WHITE, PP_ALIGN.CENTER)
        textbox(s, 2.7, y + 0.18, 9.6, 0.45, v, 18, False, INK)
    footer(s, prs, "9")

    # 10 missing
    s = new_slide(prs)
    textbox(s, 0.6, 0.4, 12, 0.5, "Bewusst nicht drin  —  mit Owner", 26, True, NAVY)
    missing = [
        ("BOM / Massen / GWP", "Fachdaten, nicht raten", "Thomas"),
        ("Sequenz, Dauer, Kosten", "AnyLogic bleibt bei Marc", "Marc"),
        ("AASX / semanticIds", "Form erst nach Interface", "Shubham + Heiko"),
        ("Vessel-Defaults final", "Katalog-Platzhalter → Marc-Override", "Marc"),
    ]
    for i, (a, b, c) in enumerate(missing):
        y = 1.2 + i * 1.25
        card(s, 0.6, y, 12.1, 1.1)
        textbox(s, 0.85, y + 0.15, 8.5, 0.4, a, 20, True, NAVY)
        textbox(s, 0.85, y + 0.55, 8.5, 0.4, b, 15, False, MUTED)
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.6), Inches(y + 0.3), Inches(2.8), Inches(0.5))
        sh.fill.solid()
        sh.fill.fore_color.rgb = ORANGE
        sh.line.fill.background()
        textbox(s, 9.6, y + 0.35, 2.8, 0.4, c, 13, True, WHITE, PP_ALIGN.CENTER)
    footer(s, prs, "10")

    # 11 then UI
    s = new_slide(prs)
    textbox(s, 0.6, 0.4, 12, 0.5, "Erst jetzt: das Working Board", 26, True, NAVY)
    textbox(
        s,
        0.6,
        1.2,
        12,
        0.8,
        "Die Tabelle ist das Ergebnis der Folien davor — nicht der ganze Stand.",
        20,
        False,
        TEAL,
    )
    bullets(
        s,
        0.7,
        2.2,
        12,
        4.2,
        [
            "Live: Assets → Waves → Economics → Vessels · Depth/Shore-Filter · Map-light",
            "Park-Dossier: Steckbrief/Site, Wetter Tag+Stunde+CSV, Akteure, Schiffseinsätze, Sim-Rollen",
            "Vessels: Typenkatalog Marc-Felder + CSV · Contracts DE light (Suche/Sort)",
            "Team-Slug lokal: research-team  ·  http://localhost:3000  ·  Branch feat/assets-ia-restructure",
        ],
        18,
    )
    footer(s, prs, "11")

    # 12 plan
    s = new_slide(prs)
    textbox(s, 0.6, 0.4, 12, 0.5, "Als Nächstes: zwei Spuren parallel", 26, True, NAVY)
    card(s, 0.6, 1.2, 6.0, 5.3)
    textbox(s, 0.85, 1.4, 5.5, 0.4, "Spur A  ·  Daten", 18, True, TEAL)
    bullets(
        s,
        0.85,
        2.0,
        5.5,
        4.2,
        [
            "Marc-IA abnehmen: Stunden-CSV, Katalog, Barge, Snapshot-Modell",
            "ERA5-Stunden AV bereits CDS (~23k h); DE-Tagesbatch optional",
            "Vessel-Typen-Katalog (8) mit Platzhaltern — Marc override",
            "Cloud nur per psql, nicht MCP-Seed",
        ],
        16,
    )
    card(s, 6.9, 1.2, 6.0, 5.3)
    textbox(s, 7.15, 1.4, 5.5, 0.4, "Spur B  ·  Interop", 18, True, TEAL)
    bullets(
        s,
        7.15,
        2.0,
        5.5,
        4.2,
        [
            "Industrie-Normen ablegen (IDTA / DPP)",
            "Mit Shubham: AAS-Schnitt Alpha Ventus",
            "Interface Agreement DB → AAS",
            "Erst dann Exporter — nicht vorher",
        ],
        16,
    )
    footer(s, prs, "12")

    # 13 asks
    s = new_slide(prs)
    textbox(s, 0.6, 0.4, 12, 0.5, "Was ich von euch brauche", 26, True, NAVY)
    asks = [
        ("Marc", "Stunden-CSV + Katalog OK? Barge eigener Typ? Sequenz-Vorlage wann?"),
        ("Thomas", "Massen / GWP kommen von dir in die DB, nicht aus der UI."),
        ("Shubham", "Normen + Mapping: du Form, ich Spalten und Keys."),
        ("Alle", "Kein Feature-Zuruf an die Tabelle. Fehlendes bleibt beim Owner."),
    ]
    for i, (who, what) in enumerate(asks):
        y = 1.2 + i * 1.25
        card(s, 0.6, y, 12.1, 1.1)
        textbox(s, 0.85, y + 0.3, 2.2, 0.5, who, 22, True, TEAL)
        textbox(s, 3.2, y + 0.32, 9.1, 0.55, what, 18, False, INK)
    footer(s, prs, "13")

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
